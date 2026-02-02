# app/services/cnc_generator.py
"""
CNC PowerPoint Generator Service

Generates complete CNC (Composizione Negoziata della Crisi) PowerPoint presentations
from backend DB data. Ported from ghigus-cnc/main.py, adapted to use DB-backed
data via cnc_data_adapter instead of Excel ranges.

All 33 slides generated (no exclusions). Data comes from cnc_data_adapter.
AI comments from cnc_ai. Company info from Case model.
"""
import io
import logging
import re
import time
from dataclasses import dataclass, field
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from sqlalchemy.orm import Session

from app.db.models.case import Case
from app.services.cnc_ai import AIResponse, generate_comment, _extract_periodo
from app.services import cnc_data_adapter as adapter

logger = logging.getLogger(__name__)

# === COLORS (same as ghigus-cnc) ===
COLORS = {
    "header_bg": RGBColor(0x1F, 0x4E, 0x79),
    "header_text": RGBColor(0xFF, 0xFF, 0xFF),
    "row_alt_bg": RGBColor(0xD6, 0xE3, 0xF8),
    "total_bg": RGBColor(0xB4, 0xC6, 0xE7),
    "title_color": RGBColor(0x1F, 0x4E, 0x79),
    "text_color": RGBColor(0x33, 0x33, 0x33),
    "footer_color": RGBColor(0x80, 0x80, 0x80),
}

MAX_WORDS_PER_SLIDE = 250


# === SLIDE DEFINITIONS ===
# All 33 slides. Each has: number, title, type, and optional data_key/prompt_type/table_ref.

@dataclass
class SlideConfig:
    numero: int
    titolo: str
    tipo: str  # titolo, sezione, tabella, commento_ai, testo_statico, testo_manuale
    data_key: str = ""  # key for cnc_data_adapter function
    prompt_type: str = ""  # for AI comment slides
    table_ref: int = 0  # slide number of the table this comment references
    note: str = ""


SLIDE_DEFINITIONS: list[SlideConfig] = [
    SlideConfig(1, "Piano di Ristrutturazione dei Debiti", "titolo"),
    SlideConfig(2, "Obiettivo e Strategia", "testo_statico"),
    SlideConfig(3, "PARTE 1 - Analisi della Situazione Attuale", "sezione"),
    SlideConfig(4, "Descrizione dell'azienda", "testo_manuale"),
    SlideConfig(5, "Analisi Economica", "tabella", data_key="analisi_economica"),
    SlideConfig(6, "Analisi Economica", "commento_ai", prompt_type="analisi_economica", table_ref=5),
    SlideConfig(7, "Analisi Finanziaria", "tabella", data_key="analisi_finanziaria"),
    SlideConfig(8, "Analisi Finanziaria", "commento_ai", prompt_type="analisi_finanziaria", table_ref=7),
    SlideConfig(9, "Stato Patrimoniale", "tabella", data_key="stato_patrimoniale"),
    SlideConfig(10, "Stato Patrimoniale", "commento_ai", prompt_type="stato_patrimoniale", table_ref=9),
    SlideConfig(11, "Posizione Finanziaria Netta", "tabella", data_key="pfn"),
    SlideConfig(12, "Posizione Finanziaria Netta", "commento_ai", prompt_type="posizione_finanziaria_netta", table_ref=11),
    SlideConfig(13, "Soluzioni per Superamento Crisi", "testo_statico", note="Introduzione ADR - contenuto standard"),
    SlideConfig(14, "Soluzioni - Piano Economico-Finanziario", "testo_statico", note="Contenuti essenziali del piano"),
    SlideConfig(15, "Soluzioni - Struttura Accordo", "testo_statico", note="Adesione creditori, percentuali richieste"),
    SlideConfig(16, "Soluzioni - Procedura Omologazione", "testo_statico", note="Iter procedurale"),
    SlideConfig(17, "Soluzioni - Monitoraggio", "testo_statico", note="Esecuzione ADR e benefici"),
    SlideConfig(18, "PARTE 2 - Il Piano", "sezione"),
    SlideConfig(19, "Rettifiche Attivo Patrimoniale", "tabella", data_key="rettifiche_attivo"),
    SlideConfig(20, "Rettifiche Attivo Patrimoniale", "commento_ai", prompt_type="rettifiche_attivo", table_ref=19),
    SlideConfig(21, "Rettifiche Passivo Patrimoniale", "tabella", data_key="rettifiche_passivo"),
    SlideConfig(22, "Rettifiche Passivo Patrimoniale", "commento_ai", prompt_type="rettifiche_passivo", table_ref=21),
    SlideConfig(23, "Tipologie Creditori - Non Aderenti", "tabella", data_key="creditori_non_aderenti"),
    SlideConfig(24, "Tipologie Creditori - Aderenti", "tabella", data_key="creditori_aderenti"),
    SlideConfig(25, "Affitto d'Azienda", "tabella", data_key="affitto"),
    SlideConfig(26, "Cessione Azienda", "tabella", data_key="cessione"),
    SlideConfig(27, "Spese di Prededuzione", "tabella", data_key="prededuzioni"),
    SlideConfig(28, "Il Piano", "sezione"),
    SlideConfig(29, "Il Piano - Flusso a Disposizione Anno 1", "tabella", data_key="flusso_anno1"),
    SlideConfig(30, "Il Piano - Flusso a Disposizione", "commento_ai", prompt_type="piano_flussi", table_ref=29),
    SlideConfig(31, "Il Piano - Soddisfacimento Creditori", "tabella", data_key="soddisfacimento_creditori"),
    SlideConfig(32, "Il Piano - Soddisfacimento Creditori", "commento_ai", prompt_type="piano_creditori", table_ref=31),
    SlideConfig(33, "Conclusioni e Sintesi", "commento_ai", prompt_type="conclusioni"),
]


def _get_slide_by_num(num: int) -> Optional[SlideConfig]:
    """Get a slide config by its number."""
    for s in SLIDE_DEFINITIONS:
        if s.numero == num:
            return s
    return None


# === DATA FETCHER ===

_DATA_KEY_MAP = {
    "analisi_economica": "get_analisi_economica",
    "analisi_finanziaria": "get_analisi_finanziaria",
    "stato_patrimoniale": "get_stato_patrimoniale",
    "pfn": "get_pfn_table",
    "rettifiche_attivo": "get_rettifiche_attivo",
    "rettifiche_passivo": "get_rettifiche_passivo",
    "creditori_non_aderenti": "get_creditori_non_aderenti",
    "creditori_aderenti": "get_creditori_aderenti",
    "affitto": "get_affitto_table",
    "cessione": "get_cessione_table",
    "prededuzioni": "get_prededuzioni_table",
    "flusso_anno1": "get_flusso_anno1",
    "soddisfacimento_creditori": "get_soddisfacimento_creditori",
}


def _fetch_table_data(
    data_key: str, db: Session, case_id: int, scenario_id: int
) -> list[list[str]]:
    """Fetch table data from cnc_data_adapter by data_key."""
    func_name = _DATA_KEY_MAP.get(data_key)
    if not func_name:
        logger.warning("Unknown data_key: %s", data_key)
        return []
    func = getattr(adapter, func_name, None)
    if func is None:
        logger.warning("Adapter function not found: %s", func_name)
        return []
    try:
        return func(db, case_id, scenario_id)
    except Exception as e:
        logger.error("Error fetching data for %s: %s", data_key, e)
        return []


# === SLIDE BUILDERS ===

def _add_title_slide(prs: Presentation, config: SlideConfig, azienda_info: dict) -> None:
    """Slide 1: Title slide with company name."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(12), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = config.titolo
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLORS["title_color"]
    p.alignment = PP_ALIGN.CENTER

    nome = azienda_info.get("nome_azienda", "Azienda")
    if nome:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4), Inches(12), Inches(0.8)
        )
        tf2 = subtitle_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = str(nome)
        p2.font.size = Pt(28)
        p2.font.color.rgb = COLORS["text_color"]
        p2.alignment = PP_ALIGN.CENTER


def _add_section_slide(prs: Presentation, config: SlideConfig) -> None:
    """Section divider slide."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3), Inches(12), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = config.titolo
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS["title_color"]
    p.alignment = PP_ALIGN.CENTER


def _add_table_slide(prs: Presentation, config: SlideConfig, table_data: list[list[str]]) -> None:
    """Data table slide with formatting."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.2), Inches(12), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = config.titolo.replace(" - Tabella", "")
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS["title_color"]

    if not table_data:
        _add_placeholder(slide, "Dati non disponibili")
        return

    # Limit table dimensions
    max_rows = min(len(table_data), 18)
    max_cols = min(len(table_data[0]) if table_data else 6, 8)
    table_data = [row[:max_cols] for row in table_data[:max_rows]]

    num_rows = len(table_data)
    num_cols = len(table_data[0])

    table = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.3), Inches(0.9),
        Inches(12.7), Inches(6.2)
    ).table

    # Column widths
    first_col_width = Inches(3.2)
    other_width = (12.7 - 3.2) / (num_cols - 1) if num_cols > 1 else 2
    table.columns[0].width = int(first_col_width)
    for i in range(1, num_cols):
        table.columns[i].width = int(Inches(other_width))

    # Populate table with formatting
    for row_idx, row_data in enumerate(table_data):
        is_header = row_idx == 0
        is_total = any("TOTALE" in str(v).upper() for v in row_data)

        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)

            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(9)
            para.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.RIGHT

            if is_header:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS["header_bg"]
                para.font.color.rgb = COLORS["header_text"]
                para.font.bold = True
            elif is_total:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS["total_bg"]
                para.font.bold = True
            elif row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS["row_alt_bg"]

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def _add_placeholder(slide, text: str) -> None:
    """Add a placeholder text."""
    box = slide.shapes.add_textbox(
        Inches(2), Inches(3), Inches(9), Inches(1)
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = COLORS["footer_color"]
    p.alignment = PP_ALIGN.CENTER


def _split_comment(text: str, max_words: int = MAX_WORDS_PER_SLIDE) -> list[str]:
    """Split text into chunks of max_words, breaking on double newlines."""
    paragraphs = text.strip().split("\n\n")
    chunks: list[str] = []
    current_chunk: list[str] = []
    current_word_count = 0

    for para in paragraphs:
        para_word_count = len(para.split())
        if current_word_count + para_word_count > max_words and current_chunk:
            chunks.append("\n\n".join(current_chunk))
            current_chunk = [para]
            current_word_count = para_word_count
        else:
            current_chunk.append(para)
            current_word_count += para_word_count

    if current_chunk:
        chunks.append("\n\n".join(current_chunk))

    return chunks if chunks else [text]


def _add_formatted_runs(paragraph, text: str, is_heading: bool, is_bullet: bool) -> None:
    """Insert text with markdown **bold** support."""
    if is_heading:
        font_size = Pt(14)
        font_color = COLORS["title_color"]
        default_bold = True
    else:
        font_size = Pt(12)
        font_color = COLORS["text_color"]
        default_bold = False

    parts = re.split(r"(\*\*.*?\*\*)", text)
    for part in parts:
        if not part:
            continue
        if part.startswith("**") and part.endswith("**"):
            run = paragraph.add_run()
            run.text = part[2:-2]
            run.font.bold = True
            run.font.size = font_size
            run.font.color.rgb = font_color
        else:
            run = paragraph.add_run()
            run.text = part
            run.font.bold = default_bold
            run.font.size = font_size
            run.font.color.rgb = font_color


def _render_comment_slide(
    prs: Presentation, title: str, comment_text: str, is_continuation: bool = False
) -> None:
    """Create a single comment slide with formatted text."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    display_title = f"{title} (continua)" if is_continuation else title
    title_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.2), Inches(12), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = display_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS["title_color"]

    # Comment box
    comment_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1), Inches(12.3), Inches(6)
    )
    tf = comment_box.text_frame
    tf.word_wrap = True

    blocks = re.split(r"\n\n+", comment_text.strip())
    first_para = True

    for block in blocks:
        lines = block.split("\n")
        for line in lines:
            line = line.strip()
            if not line:
                continue

            if first_para:
                p = tf.paragraphs[0]
                first_para = False
            else:
                p = tf.add_paragraph()

            clean_line = line.lstrip("#").strip()
            is_heading = line.startswith("#") or bool(
                re.match(r"^\d+\.\s+[A-Z]", clean_line)
            )
            is_bullet = line.startswith("- ") or line.startswith("• ")

            if is_bullet:
                clean_line = line[2:].strip()
                p.level = 1

            _add_formatted_runs(p, clean_line, is_heading, is_bullet)

            p.line_spacing = 1.15
            if is_heading:
                p.space_before = Pt(8)
                p.space_after = Pt(2)

        # Spacer after block
        spacer = tf.add_paragraph()
        spacer.space_before = Pt(4)
        spacer.font.size = Pt(4)


def _add_comment_slides(
    prs: Presentation,
    config: SlideConfig,
    table_data: list[list[str]],
    dati_sintesi: dict,
    use_mock: bool,
    stats: dict,
) -> None:
    """Generate AI comment and add slide(s), splitting if >250 words."""
    prompt_type = config.prompt_type
    periodo = _extract_periodo(table_data) if table_data else "2021-2025"

    extra_kwargs: dict = {"periodo": periodo}
    if prompt_type == "conclusioni":
        extra_kwargs = {"dati_sintesi": dati_sintesi}
    elif prompt_type in ("piano_flussi", "piano_creditori"):
        extra_kwargs = {"anno": 1}
    elif prompt_type in ("rettifiche_attivo", "rettifiche_passivo"):
        extra_kwargs = {"anno_riferimento": periodo}

    ai_start = time.time()
    try:
        response = generate_comment(
            table_data, prompt_type, use_mock=use_mock, **extra_kwargs
        )
        comment = response.content
        stats["ai_comments"] += 1
        stats["ai_tokens"] += response.tokens_used
        stats["ai_time"] += time.time() - ai_start
    except Exception as e:
        logger.error("Error generating AI comment for %s: %s", prompt_type, e)
        comment = f"[Commento AI non disponibile: {e}]"
        stats["errors"] += 1

    title = config.titolo
    chunks = _split_comment(comment)

    if len(chunks) > 1:
        logger.info(
            "Comment '%s' split across %d slides (%d words)",
            title, len(chunks), len(comment.split()),
        )

    for i, chunk in enumerate(chunks):
        _render_comment_slide(prs, title, chunk, is_continuation=(i > 0))
        stats["slides_created"] += 1


def _add_obiettivo_slide(prs: Presentation) -> None:
    """Slide 2: Obiettivo e Strategia (static content)."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    title_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.3), Inches(12), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Obiettivo e Strategia"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS["title_color"]

    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(12), Inches(5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    sections = [
        (
            "Obiettivo",
            "Garantire la continuità aziendale e soddisfare i creditori "
            "attraverso un piano di ristrutturazione sostenibile.",
        ),
        (
            "Strategia",
            "Piano di ristrutturazione dei debiti ai sensi dell'art. 57 del "
            "Codice della Crisi d'Impresa e dell'Insolvenza (CCII).",
        ),
        (
            "Base Normativa",
            "D.Lgs. 14/2019 - Codice della Crisi d'Impresa e dell'Insolvenza, "
            "con particolare riferimento agli Accordi di Ristrutturazione dei "
            "Debiti (Art. 57 e seguenti).",
        ),
    ]

    for i, (heading, body) in enumerate(sections):
        if i > 0:
            p_h = tf.add_paragraph()
        else:
            p_h = tf.paragraphs[0]
        p_h.text = heading
        p_h.font.size = Pt(20)
        p_h.font.bold = True
        p_h.font.color.rgb = COLORS["title_color"]

        p_b = tf.add_paragraph()
        p_b.text = body
        p_b.font.size = Pt(16)
        p_b.font.color.rgb = COLORS["text_color"]
        p_b.space_after = Pt(20)


def _add_descrizione_azienda_slide(prs: Presentation, azienda_info: dict) -> None:
    """Slide 4: Company description with dynamic data."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    title_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.3), Inches(12), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Descrizione dell'Azienda"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLORS["title_color"]

    nome = azienda_info.get("nome_azienda", "[Nome Azienda]")
    settore = azienda_info.get("settore", "[Settore di attività]")
    anno = azienda_info.get("anno_analisi", "2025")
    mese = azienda_info.get("mese_analisi", "N/D")

    mesi_map = {
        "GEN": "Gennaio", "FEB": "Febbraio", "MAR": "Marzo",
        "APR": "Aprile", "MAG": "Maggio", "GIU": "Giugno",
        "LUG": "Luglio", "AGO": "Agosto", "SET": "Settembre",
        "OTT": "Ottobre", "NOV": "Novembre", "DIC": "Dicembre",
    }
    mese_completo = mesi_map.get(str(mese).upper(), mese)

    if not nome or nome == "None":
        nome = "[Nome Azienda]"
    if not settore or settore == "None":
        settore = "[Settore di attività]"

    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.5), Inches(12), Inches(5)
    )
    tf = content_box.text_frame
    tf.word_wrap = True

    info_lines = [
        ("Azienda:", str(nome)),
        ("Settore:", str(settore)),
        ("Anno di analisi:", str(anno)),
        ("Periodo di riferimento:", f"Gennaio - {mese_completo} {anno}"),
    ]

    for i, (label, value) in enumerate(info_lines):
        if i > 0:
            p_label = tf.add_paragraph()
        else:
            p_label = tf.paragraphs[0]

        p_label.text = label
        p_label.font.size = Pt(18)
        p_label.font.bold = True
        p_label.font.color.rgb = COLORS["title_color"]
        p_label.space_before = Pt(15)

        p_value = tf.add_paragraph()
        p_value.text = value
        p_value.font.size = Pt(16)
        p_value.font.color.rgb = COLORS["text_color"]

    if "[" in str(nome) or "[" in str(settore):
        note_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(6.2), Inches(12), Inches(0.5)
        )
        ntf = note_box.text_frame
        np_ = ntf.paragraphs[0]
        np_.text = "Nota: completare i campi tra parentesi con i dati dell'azienda"
        np_.font.size = Pt(10)
        np_.font.italic = True
        np_.font.color.rgb = COLORS["footer_color"]


def _add_static_slide(prs: Presentation, config: SlideConfig) -> None:
    """Generic static text slide."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    title_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.2), Inches(12), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = config.titolo
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLORS["title_color"]

    if config.note:
        note_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2), Inches(12), Inches(1)
        )
        ntf = note_box.text_frame
        np_ = ntf.paragraphs[0]
        np_.text = f"[{config.note}]"
        np_.font.size = Pt(12)
        np_.font.italic = True
        np_.font.color.rgb = COLORS["footer_color"]


# === MAIN GENERATOR ===

def generate(
    db: Session,
    case_id: int,
    scenario_id: int,
    use_mock: bool = True,
) -> io.BytesIO:
    """
    Generate the complete CNC PowerPoint presentation.

    Args:
        db: SQLAlchemy session
        case_id: Case ID (integer PK or string slug used as PK)
        scenario_id: Scenario ID
        use_mock: If True, use mock AI comments; if False, call OpenAI API

    Returns:
        BytesIO containing the generated .pptx file
    """
    logger.info(
        "Generating CNC presentation for case=%s scenario=%s (mock=%s)",
        case_id, scenario_id, use_mock,
    )

    # Get company info from Case
    case = db.query(Case).filter(Case.id == str(case_id)).first()
    azienda_info: dict = {}
    if case:
        azienda_info["nome_azienda"] = case.name or "[Nome Azienda]"
        azienda_info["settore"] = case.description or "[Settore di attività]"
        azienda_info["anno_analisi"] = "2025"
        azienda_info["mese_analisi"] = "DIC"

    # Get dati_sintesi for conclusioni slide
    dati_sintesi = adapter.get_dati_sintesi(db, case_id, scenario_id)

    # Pre-fetch all table data (keyed by data_key)
    table_cache: dict[str, list[list[str]]] = {}
    for sdef in SLIDE_DEFINITIONS:
        if sdef.data_key and sdef.data_key not in table_cache:
            table_cache[sdef.data_key] = _fetch_table_data(
                sdef.data_key, db, case_id, scenario_id
            )

    # Create presentation (16:9)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    stats = {
        "slides_created": 0,
        "ai_comments": 0,
        "ai_tokens": 0,
        "ai_time": 0.0,
        "errors": 0,
    }

    # Process each slide
    for sdef in SLIDE_DEFINITIONS:
        try:
            if sdef.tipo == "titolo":
                _add_title_slide(prs, sdef, azienda_info)
                stats["slides_created"] += 1

            elif sdef.tipo == "sezione":
                _add_section_slide(prs, sdef)
                stats["slides_created"] += 1

            elif sdef.tipo == "tabella":
                table_data = table_cache.get(sdef.data_key, [])
                _add_table_slide(prs, sdef, table_data)
                stats["slides_created"] += 1

            elif sdef.tipo == "commento_ai":
                # Get table data from referenced table slide
                table_data = []
                if sdef.table_ref:
                    ref_slide = _get_slide_by_num(sdef.table_ref)
                    if ref_slide and ref_slide.data_key:
                        table_data = table_cache.get(ref_slide.data_key, [])
                _add_comment_slides(
                    prs, sdef, table_data, dati_sintesi, use_mock, stats
                )

            elif sdef.tipo in ("testo_statico", "testo_manuale"):
                if sdef.numero == 2:
                    _add_obiettivo_slide(prs)
                elif sdef.numero == 4:
                    _add_descrizione_azienda_slide(prs, azienda_info)
                else:
                    _add_static_slide(prs, sdef)
                stats["slides_created"] += 1

            else:
                logger.warning("Unknown slide type: %s", sdef.tipo)
                _add_static_slide(prs, sdef)
                stats["slides_created"] += 1

        except Exception as e:
            logger.error("Error processing slide %d (%s): %s", sdef.numero, sdef.titolo, e)
            stats["errors"] += 1

    logger.info(
        "CNC generation complete: %d slides, %d AI comments, %d tokens, %d errors",
        stats["slides_created"],
        stats["ai_comments"],
        stats["ai_tokens"],
        stats["errors"],
    )

    # Save to BytesIO
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output
